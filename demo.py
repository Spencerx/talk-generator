import argparse
import pathlib
import os.path
import random
import inflect

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from nltk.corpus import wordnet as wn
from py_thesaurus import Thesaurus
from PyDictionary import PyDictionary
from google_images_download import google_images_download

# CONSTANTS
HEIGHT = 9
WIDTH = 16
LEFTMOST = Inches(0)
TOPMOST = Inches(0)
HEIGHT_IN = Inches(HEIGHT)
WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs 
INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
CMS_TO_EMU = 360000


def get_definitions(word):
    """Get definitions of a given topic word."""
    print('******************************************')
    print('Word: {}'.format(word))
    # Get definition
    dictionary = PyDictionary()
    definitions = dictionary.meaning(word)
    if definitions:
        print('******************************************')
        print('{} word type(s)'.format(len(definitions)))
        for word_type, word_type_defs in definitions.items():
            print('******************************************')
            print('As a {}, {} can mean {} things: '.format(word_type.lower(),
                word, len(word_type_defs)))

            for word_type_def in word_type_defs:
                print('\t {}'.format(word_type_def))
        return definitions
    else:
        print('No definition found.')
        return None


def get_synonyms(word):
    """Get all synonyms for a given word."""
    print('******************************************')
    word_senses = wn.synsets(word)
    all_synonyms = []
    for ss in word_senses:
        # print(ss.name(), ss.lemma_names(), ss.definition())
        all_synonyms.extend(
            [x.lower().replace('_', ' ') for x in ss.lemma_names()])

    all_synonyms = list(set(all_synonyms))

    print('{} synonyms: '.format(len(all_synonyms)))
    for synonym in all_synonyms:
        if synonym is not word.lower():
            print('\t {}'.format(synonym))

    return all_synonyms


def get_title(synonyms):
    """Returns a template title from a source list."""
    print('******************************************')
    chosen_synonym = random.choice(synonyms)
    chosen_synonym_plural = inflect.engine().plural(chosen_synonym)
    synonym_templates = ['The Unexpected Benefits of {}',
                         'What Your Choice in {} Says About You',
                         'How to Get Rid of {}',
                         'Why {} Will Ruin Your Life',
                         'The Biggest Concerns About {}']
    chosen_template = random.choice(synonym_templates);
    return chosen_template.format(chosen_synonym_plural.title())


def get_images(synonyms, num_images):
    if num_images > 0:
        all_paths = {}
        for synonym in synonyms:
            # Get related images at 16x9 aspect ratio
            # TODO: add image filter for weird and NSFW stuff
            response = google_images_download.googleimagesdownload()
            arguments = {
                'keywords': synonym,
                'limit': num_images,
                'print_urls': True,
                'exact_size': '1600,900',
                # 'size':'large',
                # 'usage_rights':'labeled-for-noncommercial-reuse-with-modification'
            }
            # passing the arguments to the function
            paths = response.download(arguments)
            # printing absolute paths of the downloaded images
            print(paths)
            # Add to main dictionary
            all_paths[synonym] = paths[synonym]
    else:
        all_paths = None
    return all_paths


def compile_presentation(args, all_paths, title, definitions, synonyms):
    # Make a presentation
    prs = Presentation()

    # Set the height and width
    prs.slide_height = HEIGHT * INCHES_TO_EMU
    prs.slide_width = WIDTH * INCHES_TO_EMU

    # Get a default blank slide layout
    slide_layout = prs.slide_layouts[5]

    # Build an ordered list of slides for access
    slides = []

    # Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slides.append(slide)
    title_object = slide.shapes.title
    title_object.text = title
    title_object.width = WIDTH_IN
    title_object.height = HEIGHT_IN
    title_object.left = LEFTMOST
    title_object.right = TOPMOST

    slide_idx_iter = 1
    for synonym, paths in all_paths.items():
        print('***********************************')
        print('Adding slide: {}'.format(slide_idx_iter))
        slides.append(prs.slides.add_slide(slide_layout))

        # Get an image
        img_path = paths[0]

        # Add the image to the slide.
        if img_path:
            pic = slides[slide_idx_iter].shapes.add_picture(img_path, 
                LEFTMOST, TOPMOST, width=WIDTH_IN, height=HEIGHT_IN)

            # Add title to the slide
            shapes = slides[slide_idx_iter].shapes
            shapes.title.text = synonym

            # TODO: Add the text to the slide.

            slide_idx_iter += 1

    return prs, slides


def save_talk(args, prs):
    # Save the presentation
    fp = './output/' + args.topic + '-' + args.output
    # Create the parent folder if it doesn't exist
    pathlib.Path(os.path.dirname(fp)).mkdir(parents=True, exist_ok=True)
    prs.save(fp)
    print('Saved talk to {}'.format(fp))
    return True


def main(args):
    """Make a presentation with the given topic."""
    # Get definitions
    definitions = get_definitions(args)
    # Get synonyms
    synonyms = get_synonyms(args)
    # Get a title
    title = get_title(synonyms)
    # For each synonym download num_images
    all_paths = get_images(synonyms, args.num_images)
    # Compile the presentation
    prs, slides = compile_presentation(args, 
                                       all_paths=all_paths, 
                                       title=title,
                                       definitions=definitions, 
                                       synonyms=synonyms)
    if save_talk(args, prs):
        print('Successfully built talk.')


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Process some integers.')
    parser.add_argument('--output', help="Output filename.",
                        default='test.pptx', type=str)
    parser.add_argument('--topic', help="Topic of presentation.",
                        default='bagels', type=str)
    parser.add_argument('--num_images', help="Number of images per synonym.",
                        default=0, type=int)
    parser.add_argument('--num_slides', help="Number of slides to create.",
                        default=6, type=int)
    args = parser.parse_args()
    main(args)
